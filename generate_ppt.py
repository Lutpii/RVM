from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color Palette ────────────────────────────────────────────────────────────
C_DARK_BLUE   = RGBColor(0x1A, 0x37, 0x5E)   # deep navy
C_TEAL        = RGBColor(0x00, 0x89, 0x7B)   # primary accent
C_LIGHT_TEAL  = RGBColor(0xE0, 0xF2, 0xF1)   # soft teal bg
C_GREEN       = RGBColor(0x2E, 0x7D, 0x32)   # eco green
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
C_MID_GRAY    = RGBColor(0x90, 0xA4, 0xAE)
C_TEXT        = RGBColor(0x21, 0x21, 0x21)
C_SUBTITLE    = RGBColor(0x45, 0x55, 0x68)
C_BORDER      = RGBColor(0xB2, 0xDF, 0xDB)
C_ORANGE      = RGBColor(0xFF, 0x8F, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # truly blank layout

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_size=14, bold=False, color=C_TEXT,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def section_header(slide, title, subtitle=None):
    """Colored top bar with section title."""
    add_rect(slide, 0, 0, 13.33, 1.1, fill=C_DARK_BLUE)
    add_text_box(slide, title,
                 0.35, 0.12, 10, 0.65,
                 font_size=26, bold=True, color=C_WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     0.35, 0.72, 10, 0.35,
                     font_size=13, color=C_TEAL)
    # thin teal underline
    add_rect(slide, 0, 1.1, 13.33, 0.06, fill=C_TEAL)


def bullet_box(slide, items, l, t, w, h,
               font_size=13, color=C_TEXT, spacing=1.4):
    """Render a bullet list into a text box."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(spacing)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def card(slide, l, t, w, h, title, body_lines,
         title_color=C_DARK_BLUE, title_size=13,
         body_size=11.5, icon=None):
    """Rounded card with title + bullets."""
    add_rect(slide, l, t, w, h, fill=C_LIGHT_GRAY, line=C_BORDER, line_w=Pt(1))
    # top accent bar
    add_rect(slide, l, t, w, 0.07, fill=C_TEAL)
    # icon + title
    header = (icon + "  " + title) if icon else title
    add_text_box(slide, header, l + 0.15, t + 0.12, w - 0.25, 0.35,
                 font_size=title_size, bold=True, color=title_color)
    bullet_box(slide, body_lines, l + 0.15, t + 0.48, w - 0.25, h - 0.55,
               font_size=body_size)


def flow_arrow(slide, x, y, length=0.55):
    """Draw right-pointing arrow shape."""
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        13,  # right arrow
        Inches(x), Inches(y), Inches(length), Inches(0.35)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_TEAL
    shape.line.fill.background()
    return shape


def flow_box(slide, l, t, w, h, text, bg=C_DARK_BLUE, fg=C_WHITE, size=10.5):
    add_rect(slide, l, t, w, h, fill=bg)
    add_text_box(slide, text, l + 0.05, t + 0.06, w - 0.1, h - 0.1,
                 font_size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)

# full dark-blue background
add_rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK_BLUE)

# green accent strip (left)
add_rect(slide, 0, 0, 0.45, 7.5, fill=C_TEAL)

# decorative teal circle (top-right)
circ = slide.shapes.add_shape(9,  # ellipse
    Inches(10.5), Inches(-1.2), Inches(4.5), Inches(4.5))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x00, 0x6E, 0x65)
circ.line.fill.background()

circ2 = slide.shapes.add_shape(9,
    Inches(11.2), Inches(5.2), Inches(3.2), Inches(3.2))
circ2.fill.solid(); circ2.fill.fore_color.rgb = RGBColor(0x00, 0x59, 0x51)
circ2.line.fill.background()

# RVM icon placeholder (recycling symbol as text)
add_text_box(slide, "♻", 0.9, 0.6, 2, 1.6,
             font_size=72, bold=True, color=C_TEAL, align=PP_ALIGN.LEFT)

# Main title
add_text_box(slide, "Reverse Vending Machine",
             0.9, 1.9, 9.5, 1.0,
             font_size=38, bold=True, color=C_WHITE)
add_text_box(slide, "(RVM) Smart Recycling System",
             0.9, 2.75, 9.5, 0.8,
             font_size=30, bold=False, color=C_TEAL)

# Divider
add_rect(slide, 0.9, 3.6, 6.0, 0.05, fill=C_TEAL)

# Subtitle
add_text_box(slide, "Collaboration Presentation",
             0.9, 3.75, 9, 0.5,
             font_size=15, color=C_MID_GRAY)

# From → To
add_text_box(slide, "UMPSA  →  DSME ENGINEERING",
             0.9, 4.25, 9, 0.55,
             font_size=18, bold=True, color=C_WHITE)

# Presenter & date
add_text_box(slide, "Presented by:  Adi Luthfi Nur Roki",
             0.9, 5.1, 7, 0.4,
             font_size=12.5, color=C_MID_GRAY)
add_text_box(slide, "Universiti Malaysia Pahang Al-Sultan Abdullah (UMPSA)  |  2026",
             0.9, 5.52, 9, 0.35,
             font_size=11, italic=True, color=C_MID_GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Background / Introduction
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
section_header(slide, "Background & Introduction",
               "Why a Reverse Vending Machine?")

# Problem column
add_rect(slide, 0.35, 1.35, 5.85, 0.42, fill=C_DARK_BLUE)
add_text_box(slide, "  The Problem", 0.35, 1.35, 5.85, 0.42,
             font_size=14, bold=True, color=C_WHITE)

bullet_box(slide, [
    "●  Recycling rates remain critically low in Malaysia — only ~23% of solid waste is recycled",
    "●  Lack of incentive drives low public participation in recycling programmes",
    "●  Manual sorting is error-prone, inefficient, and labour-intensive",
    "●  Recyclable waste often ends up in landfills due to poor collection infrastructure",
    "●  Communities lack real-time visibility into recycling bin capacity",
], 0.45, 1.85, 5.65, 3.8, font_size=12.5, spacing=3.5)

# Solution column
add_rect(slide, 6.55, 1.35, 6.42, 0.42, fill=C_TEAL)
add_text_box(slide, "  Our Solution: RVM System", 6.55, 1.35, 6.42, 0.42,
             font_size=14, bold=True, color=C_WHITE)

bullet_box(slide, [
    "●  Smart kiosk accepts recyclable materials: Aluminum, Plastic, Glass & Paper",
    "●  AI-powered material classification using YOLOv8n deep learning model",
    "●  Gamified reward system — users earn points redeemable for incentives",
    "●  Real-time bin level monitoring and machine status tracking",
    "●  Full-stack web platform: User app + Admin dashboard + AI service",
], 6.65, 1.85, 6.2, 3.8, font_size=12.5, spacing=3.5)

# Context banner
add_rect(slide, 0.35, 5.9, 12.62, 1.15, fill=C_LIGHT_TEAL, line=C_BORDER, line_w=Pt(1))
add_text_box(slide,
    "This project is developed by UMPSA students as a collaboration deliverable for DSME Engineering. "
    "The system bridges sustainable recycling practices with modern IoT and AI technologies, "
    "creating a scalable solution for smart waste management.",
    0.55, 5.98, 12.2, 1.0, font_size=12, italic=True, color=C_SUBTITLE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Dashboard Overview
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
section_header(slide, "Dashboard Overview",
               "Unified interface for users, admins, and kiosk operators")

# User Dashboard card
add_rect(slide, 0.35, 1.35, 5.85, 5.7, fill=C_LIGHT_GRAY, line=C_BORDER, line_w=Pt(1))
add_rect(slide, 0.35, 1.35, 5.85, 0.08, fill=C_TEAL)
add_text_box(slide, "👤  User Dashboard", 0.5, 1.5, 5.5, 0.45,
             font_size=15, bold=True, color=C_DARK_BLUE)

user_items = [
    "●  Total Points Balance — live tracking of earned rewards",
    "●  Recycling History — session logs with material breakdown",
    "●  Points History — detailed ledger of earned & deducted points",
    "●  Active Session Status — real-time recycling session tracker",
    "●  Profile Management — update name, language, & theme",
    "●  QR Code Scanner — link device to RVM kiosk for session start",
    "●  Bilingual UI — toggle between English and Bahasa Indonesia",
    "●  Dark / Light theme switch for comfortable usage",
]
bullet_box(slide, user_items, 0.5, 2.05, 5.6, 4.7, font_size=12, spacing=3)

# Admin Dashboard card
add_rect(slide, 6.55, 1.35, 6.42, 5.7, fill=C_LIGHT_GRAY, line=C_BORDER, line_w=Pt(1))
add_rect(slide, 6.55, 1.35, 6.42, 0.08, fill=C_DARK_BLUE)
add_text_box(slide, "🛠  Admin Dashboard", 6.7, 1.5, 6.1, 0.45,
             font_size=15, bold=True, color=C_DARK_BLUE)

admin_items = [
    "●  System Statistics — total users, sessions, items recycled & points issued",
    "●  User Management — view, edit, and remove user accounts",
    "●  Machine Management — add/edit RVM kiosks with GPS location",
    "●  Bin Level Monitoring — track fill-level per material type per machine",
    "●  Transaction Logs — full audit trail of all recycling transactions",
    "●  Reward Config — configure points per material type",
    "●  CSV Export — download data for reporting and analysis",
    "●  Admin Action Log — accountability trail for all admin operations",
]
bullet_box(slide, admin_items, 6.7, 2.05, 6.15, 4.7, font_size=12, spacing=3)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — System Workflow (Business Flow)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
section_header(slide, "System Workflow",
               "End-to-end recycling session business flow")

# Phase labels
def phase_label(slide, text, l, t, color=C_DARK_BLUE):
    add_text_box(slide, text, l, t, 2.8, 0.32,
                 font_size=10.5, bold=True, color=color, italic=True)

# Row 1: User side
phase_label(slide, "Phase 1 — User Onboarding", 0.35, 1.22, C_DARK_BLUE)

steps_r1 = [
    ("Register /\nLogin", 0.35),
    ("Scan QR\nCode", 2.25),
    ("Select\nMaterial", 4.15),
    ("Open Kiosk\nLid", 6.05),
    ("Insert\nItem", 7.95),
]
arw_positions_r1 = [1.52, 3.42, 5.32, 7.22]

for label, x in steps_r1:
    flow_box(slide, x, 1.6, 1.58, 0.75, label, bg=C_DARK_BLUE, fg=C_WHITE, size=10)
for x in arw_positions_r1:
    flow_arrow(slide, x, 1.72, 0.55)

# Extra arrow + final box in row 1
flow_arrow(slide, 9.85, 1.72, 0.55)
flow_box(slide, 10.45, 1.6, 1.58, 0.75, "Conveyor\nProcess", bg=C_DARK_BLUE, fg=C_WHITE, size=10)
flow_arrow(slide, 12.08, 1.72, 0.55)
flow_box(slide, 12.68, 1.6, 0.3, 0.75, "↓", bg=C_TEAL, fg=C_WHITE, size=10)

# Row 2: AI & System side
phase_label(slide, "Phase 2 — AI Classification & Validation", 0.35, 2.6, C_TEAL)

steps_r2 = [
    ("Capture\nImage", 12.38),
    ("AI Classify\n(YOLOv8n)", 10.48),
    ("Confidence\nCheck ≥60%", 8.58),
    ("Weigh\nItem", 6.68),
    ("Calculate\nPoints", 4.78),
]
for label, x in steps_r2:
    flow_box(slide, x, 3.0, 1.58, 0.75, label, bg=C_TEAL, fg=C_WHITE, size=10)

# arrows (going right-to-left visually, drawn left-to-right)
arw_r2 = [11.32, 9.42, 7.52, 5.62]
for x in arw_r2:
    shape = slide.shapes.add_shape(14,  # left arrow
        Inches(x), Inches(3.12), Inches(0.55), Inches(0.35))
    shape.fill.solid(); shape.fill.fore_color.rgb = C_TEAL
    shape.line.fill.background()

# small up arrow connecting rows
flow_box(slide, 12.68, 2.48, 0.3, 0.55, "↓", bg=C_TEAL, fg=C_WHITE, size=10)

# Row 3: Completion
phase_label(slide, "Phase 3 — Session Completion", 0.35, 4.0, C_GREEN)

steps_r3 = [
    ("Update Bin\nLevel", 0.35),
    ("Save\nTransaction", 2.25),
    ("Update\nUser Points", 4.15),
    ("Session\nSummary", 6.05),
    ("End\nSession", 7.95),
]
arw_positions_r3 = [1.52, 3.42, 5.32, 7.22]

for label, x in steps_r3:
    flow_box(slide, x, 4.38, 1.58, 0.75, label,
             bg=C_GREEN, fg=C_WHITE, size=10)
for x in arw_positions_r3:
    flow_arrow(slide, x, 4.5, 0.55)

# connect row2 end to row3 start
flow_box(slide, 3.58, 3.0, 0.3, 0.75, "↓", bg=C_TEAL, fg=C_WHITE, size=10)
flow_box(slide, 3.58, 3.78, 0.3, 0.62, "↓", bg=C_GREEN, fg=C_WHITE, size=10)

# Reject path note
add_rect(slide, 9.7, 3.78, 3.27, 0.58, fill=RGBColor(0xFF, 0xEB, 0xEE), line=RGBColor(0xEF, 0x9A, 0x9A), line_w=Pt(1))
add_text_box(slide, "✗  If AI confidence < 60% or mismatch:\n    Item rejected → −10 points deducted",
             9.85, 3.82, 3.1, 0.5, font_size=10, color=RGBColor(0xC6, 0x28, 0x28))

# Legend
add_rect(slide, 0.35, 5.48, 12.62, 0.55, fill=C_LIGHT_TEAL, line=C_BORDER, line_w=Pt(1))
add_text_box(slide,
    "■ Navy = User Actions     ■ Teal = AI/System Processing     ■ Green = Completion Steps     ■ Red = Rejection Path",
    0.5, 5.54, 12.3, 0.4, font_size=11, color=C_SUBTITLE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Main Features
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
section_header(slide, "Main Features",
               "Core capabilities of the RVM Smart Recycling System")

features = [
    ("🤖  AI Material Classification",
     ["YOLOv8n deep learning model", "4 classes: Aluminum, Plastic, Glass, Paper",
      "60% confidence threshold", "Bounding box detection & scoring", "Mock fallback mode for testing"],
     0.35, 1.35, 3.9, 3.05),

    ("🔐  Multi-Auth System",
     ["Email & password registration", "WhatsApp OTP via Fonnte API",
      "Google OAuth 2.0 sign-in", "Laravel Sanctum token auth", "Role-based access: User / Admin"],
     4.55, 1.35, 3.9, 3.05),

    ("📊  Real-Time Monitoring",
     ["Live bin-level per material type", "Machine status: Active / Maintenance",
      "GPS location with Google Maps", "Admin alerts on near-full bins", "Session & transaction audit logs"],
     8.75, 1.35, 4.2, 3.05),

    ("🏆  Points & Rewards Engine",
     ["Earn points per gram of material", "Aluminum 10 pts | Plastic 8 pts",
      "Glass 6 pts | Paper 5 pts /100g", "−10 pts for rejected items", "Full points history ledger"],
     0.35, 4.6, 3.9, 2.55),

    ("📱  User Experience",
     ["14-step guided recycling session", "QR code kiosk linking",
      "Bilingual: English & Bahasa Indonesia", "Dark / Light theme toggle", "Mobile-responsive PWA-ready UI"],
     4.55, 4.6, 3.9, 2.55),

    ("⚙  Tech Stack",
     ["Frontend: Vue.js 3 + Pinia + Vite", "Backend: Laravel 10 + MySQL",
      "AI Service: Python Flask + YOLOv8n", "Auth: Sanctum + Socialite + Fonnte",
      "Charts: Chart.js  |  QR: qrcode.js"],
     8.75, 4.6, 4.2, 2.55),
]

for title, items, l, t, w, h in features:
    card(slide, l, t, w, h, title, items, title_size=12.5, body_size=11)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Expected Output
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
section_header(slide, "Expected Output",
               "Measurable deliverables and outcomes of the RVM System")

# KPI boxes (top row)
kpis = [
    ("Higher\nRecycling Rate", "Incentive-driven engagement boosts public participation through the points-reward mechanism"),
    ("Accurate AI\nClassification", "YOLOv8n model achieves reliable material identification, reducing human sorting errors"),
    ("Data-Driven\nInsights", "Admin dashboard provides actionable analytics for waste management decision-making"),
    ("Scalable\nInfrastructure", "Multi-machine support with centralised monitoring enables city-wide deployment"),
]
kw = 2.95
for i, (title, desc) in enumerate(kpis):
    x = 0.35 + i * (kw + 0.18)
    add_rect(slide, x, 1.3, kw, 1.65, fill=C_DARK_BLUE)
    add_text_box(slide, title, x + 0.12, 1.38, kw - 0.2, 0.65,
                 font_size=13, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x + 0.1, 1.98, kw - 0.18, 0.95,
                 font_size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

# Detailed deliverables
add_rect(slide, 0.35, 3.15, 12.62, 0.38, fill=C_TEAL)
add_text_box(slide, "  System Deliverables", 0.45, 3.18, 8, 0.35,
             font_size=13, bold=True, color=C_WHITE)

deliverables_left = [
    "✓  Fully functional full-stack web application (Frontend + Backend + AI service)",
    "✓  Admin panel with user, machine & transaction management",
    "✓  AI classification service with YOLOv8n model (best.pt)",
    "✓  QR-based kiosk linking for seamless session initiation",
    "✓  Bilingual interface (English & Bahasa Indonesia) with theme support",
]
deliverables_right = [
    "✓  Points system with earn/deduct logic and complete audit trail",
    "✓  WhatsApp OTP and Google OAuth multi-authentication",
    "✓  Real-time bin monitoring with GPS machine locations",
    "✓  CSV data export for reporting and external integration",
    "✓  Documented API endpoints for future IoT hardware integration",
]

bullet_box(slide, deliverables_left,  0.45, 3.6, 6.15, 3.5, font_size=12, spacing=3)
bullet_box(slide, deliverables_right, 6.7,  3.6, 6.25, 3.5, font_size=12, spacing=3)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Conclusion
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=C_WHITE)
section_header(slide, "Conclusion",
               "Summary & path forward for UMPSA × DSME Engineering collaboration")

# Summary box
add_rect(slide, 0.35, 1.35, 12.62, 1.55, fill=C_LIGHT_TEAL, line=C_BORDER, line_w=Pt(1))
add_text_box(slide,
    "The RVM Smart Recycling System demonstrates a complete, production-ready solution that combines "
    "AI-powered material classification, gamified incentives, and real-time monitoring to drive "
    "sustainable recycling behaviour. Developed at UMPSA and presented to DSME Engineering as a "
    "collaboration deliverable, the system is designed for scalability, real-world deployment, and "
    "future integration with physical RVM hardware.",
    0.55, 1.42, 12.2, 1.42, font_size=13, color=C_TEXT)

# Three columns: Achievements | Next Steps | Collaboration
col_w = 3.9
titles  = ["Achievements", "Next Steps", "Collaboration Opportunities"]
col_fill = [C_DARK_BLUE, C_TEAL, C_GREEN]
col_items = [
    ["✓  Full-stack web system delivered",
     "✓  AI classification integrated (YOLOv8n)",
     "✓  14-step recycling session flow",
     "✓  Admin & user dashboards live",
     "✓  Multi-auth & bilingual support",
     "✓  Points reward engine functional"],
    ["→  Connect to physical RVM hardware via API",
     "→  Deploy to cloud (AWS / Azure / GCP)",
     "→  Mobile app (React Native / Flutter)",
     "→  Expand AI to detect more material types",
     "→  Integrate redemption / cashout system",
     "→  Conduct pilot testing at UMPSA campus"],
    ["⟷  Provide hardware kiosk specs to UMPSA team",
     "⟷  Co-develop IoT sensor integration layer",
     "⟷  Joint testing & field validation at DSME sites",
     "⟷  Share domain expertise on waste management",
     "⟷  Scale deployment to DSME partner locations",
     "⟷  Explore research publication opportunities"],
]

for i, (title, items, fill) in enumerate(zip(titles, col_items, col_fill)):
    x = 0.35 + i * (col_w + 0.22)
    add_rect(slide, x, 3.1, col_w, 0.42, fill=fill)
    add_text_box(slide, title, x + 0.1, 3.15, col_w - 0.15, 0.35,
                 font_size=13, bold=True, color=C_WHITE)
    bullet_box(slide, items, x + 0.1, 3.6, col_w - 0.15, 2.9, font_size=11.5, spacing=3)

# Thank you footer
add_rect(slide, 0, 6.75, 13.33, 0.75, fill=C_DARK_BLUE)
add_text_box(slide,
    "Thank you  |  Adi Luthfi Nur Roki  |  UMPSA  →  DSME Engineering  |  2026",
    0, 6.82, 13.33, 0.5,
    font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
output_path = r"D:\UMPSA\Semester 2\Web\RVM\RVM_System_Presentation_UMPSA_DSME.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
